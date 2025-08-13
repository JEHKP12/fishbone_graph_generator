import math
import pandas as pd
import matplotlib.pyplot as plt
from matplotlib.patches import Polygon, Wedge
from pptx import Presentation
from pptx.util import Inches
import streamlit as st
import io

# ─────────────────────────
# CONFIGURATION
# ─────────────────────────
FISH_COLOR = (245/255, 130/255, 31/255)      
BONE_SPACING = 4.0                        # distance between category "bones"

# ──────────────────────────────────────────────────────────────────────────────────────────────
# 0) Adjust the length of the fishbone tail based on the number of categories
# ──────────────────────────────────────────────────────────────────────────────────────────────
def tail_offset(n_cat: int) -> float:
    if n_cat < 3:
        return -2
    elif n_cat < 5:
        return -4
    else:
        return -5.5

# ─────────────────────────
# 3) HELPER FUNCTIONS
# ─────────────────────────
def problems(ax, label, p_x, p_y, angle_y, bone_length=180, scale=0.55):
    FIXED_ANGLE_DEGREES = 52 if angle_y > 0 else -52
    ARROW_LENGTH = bone_length
    
    angle_rad = math.radians(FIXED_ANGLE_DEGREES)
    x_off = -ARROW_LENGTH * math.cos(angle_rad)
    y_off = -ARROW_LENGTH * math.sin(angle_rad)
    
    ax.annotate(label.upper(), xy=(p_x, p_y), xytext=(x_off, y_off),
                fontsize=8, weight='bold', color='white',
                xycoords='data', textcoords='offset points',
                ha='center', va='center',
                arrowprops=dict(arrowstyle='->', facecolor='black'),
                bbox=dict(boxstyle='square', facecolor=FISH_COLOR, pad=0.8))

def causes(ax, items, c_x, c_y, top=True, scale=0.45):
    offsets = [[0.02, 0], [0.23, 0.5], [-0.46, -1],
               [0.69, 1.5], [-0.92, -2], [1.15, 2.5]]
    FIXED_ARROW_LENGTH = -20.0 
    # Ensure items is a list of strings
    items = [str(item) for item in items if pd.notna(item) and str(item).strip() != '']
    # Limit items to the number of available offsets
    items = items[:len(offsets)]
    for idx, txt in enumerate(items):
        c_x -= offsets[idx][0]
        c_y += offsets[idx][1] if top else -offsets[idx][1]
        ax.annotate(txt, xy=(c_x, c_y), 
                    xytext=(FIXED_ARROW_LENGTH, -0.3),
                    fontsize=6, ha='right', va="center",
                    xycoords='data', textcoords='offset points',
                    arrowprops=dict(arrowstyle='->', facecolor='black'))

def draw_spine(ax, xmin, xmax, head_radius, main_problem, problem_fontsize):
    ax.plot([xmin-0.1, xmax], [0, 0], color=FISH_COLOR, linewidth=2)
    ax.add_patch(Wedge((xmax, 0), head_radius, 270, 90, fc=FISH_COLOR))
    ax.text(xmax + head_radius / 6, -0.05, main_problem.upper(),
            fontsize=problem_fontsize, weight='bold', color='white')
    tail = [[xmin-0.8, 0.8], [xmin-0.8, -0.8], [xmin, 0]]
    ax.add_patch(Polygon(tail, fc=FISH_COLOR))

def draw_body(ax, categories, head_radius, problem_fontsize, main_problem, bone_length=180):
    n_cat = len(categories)
    length = math.ceil(n_cat / 2) - 1
    x_tail = tail_offset(n_cat) - length
    x_head = 2 + length

    draw_spine(ax, x_tail, x_head, head_radius, main_problem, problem_fontsize)

    offset = 0
    anchor_xy = [1.55, 0.5]

    for idx, (cat, sub) in enumerate(categories.items()):
        upper = idx % 2 == 0
        cause_y = 1.7 if upper else -1.7
        angle_y = 16 if upper else -16

        p_x = anchor_xy[0] + length + offset
        c_x = anchor_xy[1] + length + offset
        if not upper:
            offset -= BONE_SPACING

        problems(ax, cat, p_x, 0, angle_y, bone_length)
        if sub:  # Only call causes if sub is not empty
            causes(ax, sub, c_x, cause_y, top=upper)

# Streamlit App
st.title("Fishbone Diagram Generator")

# File Uploader
uploaded_file = st.file_uploader("Upload Excel file (fishbone.xlsx format)", type=["xlsx"])

# Input Parameters
head_radius = st.number_input("Fish Head Radius", value=1.5, min_value=0.1, step=0.1)
problem_fontsize = st.number_input("Main Problem Font Size", value=9.0, min_value=1.0, step=1.0)
bone_length = st.number_input("Bone Length", value=180.0, min_value=10.0, step=10.0)

if uploaded_file is not None:
    if st.button("Generate Preview"):
        try:
            df = pd.read_excel(uploaded_file, engine='openpyxl')
            if 'Main Problem' not in df.columns:
                st.error("Column 'Main Problem' is not found in Excel. Please ensure the Excel file has a 'Main Problem' column.")
                st.stop()
            main_problem = df['Main Problem'].dropna().iloc[0]
            if pd.isna(main_problem) or str(main_problem).strip() == '':
                st.error("Main Problem is empty or invalid. Please provide a valid main problem in the Excel file.")
                st.stop()
            categories = {}
            for col in df.columns:
                if col == 'Main Problem':
                    continue
                causes = df[col].dropna().tolist()
                if causes:  # Only include non-empty cause lists
                    categories[col] = causes
            if not categories:
                st.error("No valid categories found in the Excel file. Please ensure there are columns with causes.")
                st.stop()

            # Calculate dynamic limits
            n_cat = len(categories)
            length = math.ceil(n_cat / 2) - 1
            x_tail = tail_offset(n_cat) - length
            x_head = 2 + length

            fig, ax = plt.subplots(figsize=(10, 5.625))
            ax.set_xlim(x_tail - 2, x_head + head_radius + 2)
            ax.set_ylim(-6, 6)
            ax.axis('off')

            draw_body(ax, categories, head_radius, problem_fontsize, main_problem, bone_length)

            st.pyplot(fig)

            # Download PNG
            png_output = io.BytesIO()
            fig.savefig(png_output, format="png")
            png_output.seek(0)
            st.download_button(
                label="Download as PNG",
                data=png_output,
                file_name="fishbone.png",
                mime="image/png"
            )

            # Download PPTX
            pptx_io = io.BytesIO()
            temp_png_io = io.BytesIO()
            fig.savefig(temp_png_io, format="png")
            temp_png_io.seek(0)

            prs = Presentation()
            blank_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(blank_layout)
            left = Inches(0.5)
            top = Inches(0.5)
            slide.shapes.add_picture(temp_png_io, left, top, height=Inches(4))
            prs.save(pptx_io)
            pptx_io.seek(0)

            st.download_button(
                label="Download as PPTX",
                data=pptx_io,
                file_name="fishbone.pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
            )
        except Exception as e:
            st.error(f"Error generating diagram: {str(e)}")
            st.stop()
else:
    st.info("Please upload an Excel file to generate the diagram.")
